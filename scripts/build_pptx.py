# -*- coding: utf-8 -*-
"""
Assemble MOBILE/Vitech-Group-Presentation.pptx from the slide images captured by
scripts/capture_slides_for_pptx.mjs.

Each web slide becomes one landscape 16:9 PowerPoint slide holding a single
full-bleed picture (pixel-identical to the web deck). The captured JPEG for a
video slide already has a real video frame AND the title baked in, so a plain
picture reproduces it exactly.

DEFAULT = static deck (no embedded video). This is what mobile users need: the
iOS Files/attachment preview (Quick Look) cannot play embedded video and
mis-renders slides that stack a movie + poster + text overlay (heading "melts"/
doubles, some frames go blank white). Making every slide a single picture makes
the 7 former-video slides identical in structure to the normal slides that
already render perfectly everywhere, and shrinks the file ~278 MB -> ~45 MB.

Pass --with-video to instead embed the original .mp4s in place (heavier file,
tap-to-play in the Microsoft PowerPoint app). Even then we do NOT add the
transparent-titlebar overlay that caused the Quick Look doubling; the title
stays baked into the poster.

Prereqs: python-pptx, Pillow. Run after the capture step:
    python3 scripts/build_pptx.py                 # static, phone-friendly (default)
    python3 scripts/build_pptx.py --with-video    # embed videos (PowerPoint-app file)
"""
import argparse
import json
import os
import sys

from PIL import Image
from pptx import Presentation
from pptx.util import Emu

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
BUILD = os.path.join(ROOT, "MOBILE", "build")
SLIDES_DIR = os.path.join(BUILD, "slides")
OVERLAYS_DIR = os.path.join(BUILD, "overlays")
POSTERS_DIR = os.path.join(BUILD, "posters")
OVERLAYS = os.path.join(BUILD, "video_overlays.json")
OUT_PPTX = os.path.join(ROOT, "MOBILE", "Vitech-Group-Presentation.pptx")

# The capture viewport is 1920x1080; a 16:9 widescreen slide is 13.333in x 7.5in.
# 12192000 EMU / 1920 px == 6858000 EMU / 1080 px == exactly 6350 EMU per pixel.
STAGE_W, STAGE_H = 1920, 1080
SLIDE_W_EMU, SLIDE_H_EMU = 12192000, 6858000
EMU_PER_PX = 6350


def emu(px):
    return Emu(int(round(px * EMU_PER_PX)))


def clamp_rect(x, y, w, h):
    x = max(0, min(x, STAGE_W))
    y = max(0, min(y, STAGE_H))
    w = max(1, min(w, STAGE_W - x))
    h = max(1, min(h, STAGE_H - y))
    return x, y, w, h


def crop(slide_jpg, rect, out_path):
    x, y, w, h = rect
    with Image.open(slide_jpg) as im:
        im.convert("RGB").crop((x, y, x + w, y + h)).save(out_path, "JPEG", quality=92)
    return out_path


def main():
    ap = argparse.ArgumentParser(description="Build the mobile Vitech PPTX.")
    ap.add_argument(
        "--with-video", action="store_true",
        help="Embed the .mp4 videos in place (heavier file, tap-to-play in the PowerPoint "
             "app). Default is a static still-image deck that renders cleanly everywhere "
             "including the iOS Files preview.",
    )
    ap.add_argument("--out", default=OUT_PPTX, help="Output .pptx path.")
    args = ap.parse_args()

    with open(OVERLAYS) as f:
        data = json.load(f)
    order = data["slides"]           # capture order == DOM/display order 1..N
    overlays = data.get("overlays", {})

    if args.with_video:
        os.makedirs(POSTERS_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Emu(SLIDE_W_EMU)   # 13.333in  -> horizontal / landscape
    prs.slide_height = Emu(SLIDE_H_EMU)  # 7.5in
    blank = prs.slide_layouts[6]

    movie_count = 0
    for n in order:
        jpg = os.path.join(SLIDES_DIR, "slide-%03d.jpg" % n)
        if not os.path.exists(jpg):
            print("  ! missing %s" % jpg, file=sys.stderr)
            continue

        slide = prs.slides.add_slide(blank)
        # Full-bleed slide image. In static mode this IS the slide (the captured JPEG
        # already holds the video frame + title), so every slide is one clean picture.
        slide.shapes.add_picture(jpg, Emu(0), Emu(0), width=Emu(SLIDE_W_EMU), height=Emu(SLIDE_H_EMU))

        if not args.with_video:
            continue

        ov = overlays.get(str(n))
        if not ov:
            continue

        # --with-video only: embed each foreground video over its exact rectangle,
        # with a single poster cropped from the slide image (no transparent-titlebar
        # overlay -- that 3rd layer is what iOS Quick Look doubled/smeared).
        for i, v in enumerate(ov["videos"]):
            x, y, w, h = clamp_rect(v["x"], v["y"], v["w"], v["h"])
            movie_path = os.path.join(ROOT, v["file"])
            if not os.path.exists(movie_path):
                print("  ! missing video %s" % movie_path, file=sys.stderr)
                continue
            poster = crop(jpg, (x, y, w, h), os.path.join(POSTERS_DIR, "poster-%03d-%d.jpg" % (n, i)))
            slide.shapes.add_movie(
                movie_path, emu(x), emu(y), emu(w), emu(h),
                poster_frame_image=poster, mime_type="video/mp4",
            )
            movie_count += 1

        # Re-stamp opaque text panels (cards) that overlap the video, on top.
        for j, c in enumerate(ov.get("chrome", [])):
            x, y, w, h = clamp_rect(c["x"], c["y"], c["w"], c["h"])
            strip = crop(jpg, (x, y, w, h), os.path.join(POSTERS_DIR, "chrome-%03d-%d.jpg" % (n, j)))
            slide.shapes.add_picture(strip, emu(x), emu(y), width=emu(w), height=emu(h))

    out = os.path.abspath(args.out)
    os.makedirs(os.path.dirname(out), exist_ok=True)
    prs.save(out)

    size_mb = os.path.getsize(out) / (1024 * 1024)
    mode = "with embedded video" if args.with_video else "static (still frames, phone-friendly)"
    print("Wrote %s" % out)
    print("  mode: %s" % mode)
    print("  slides: %d   embedded videos: %d   size: %.1f MB" % (len(prs.slides._sldIdLst), movie_count, size_mb))


if __name__ == "__main__":
    main()
